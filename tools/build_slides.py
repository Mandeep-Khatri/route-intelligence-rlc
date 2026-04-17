#!/usr/bin/env python3
"""Build RLC hackathon PowerPoint. Run from repo root: python tools/build_slides.py"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "figures"
OUT = ROOT / "RLC_Hackathon_Presentation.pptx"

TEAL = RGBColor(13, 148, 136)
TEAL_DARK = RGBColor(15, 118, 110)
SLATE = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)


def _style_title(shape) -> None:
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEAL_DARK


def _style_body(shape) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    for i, p in enumerate(tf.paragraphs):
        p.font.size = Pt(20) if i == 0 else Pt(18)
        p.font.color.rgb = SLATE if i == 0 else MUTED
        p.space_after = Pt(10)
        if i == 0:
            p.font.bold = True


def add_bullet_slide(prs, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _style_title(slide.shapes.title)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(17)
        p.font.color.rgb = SLATE
        p.space_after = Pt(8)


def add_image_slide(prs, title: str, image_path: Path, caption: str = "") -> None:
    layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    _style_title(slide.shapes.title)

    if not image_path.is_file():
        return

    # Place image below title
    left = Inches(0.6)
    top = Inches(1.35)
    width = Inches(9.0)
    slide.shapes.add_picture(str(image_path), left, top, width=width)

    if caption:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(9), Inches(0.6))
        tf = box.text_frame
        tf.paragraphs[0].text = caption
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = MUTED
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Rescuing Leftover Cuisine"
    sub = slide.placeholders[1]
    sub.text = (
        "Data-driven impact for food rescue\n"
        "JPMorgan Chase Hackathon · NYC · Boston · Chicago · Los Angeles"
    )
    _style_title(slide.shapes.title)
    for p in sub.text_frame.paragraphs:
        p.font.size = Pt(20)
        p.font.color.rgb = MUTED
        p.alignment = PP_ALIGN.CENTER

    add_bullet_slide(
        prs,
        "The problem we addressed",
        [
            "Surplus food from businesses often becomes waste instead of meals.",
            "RLC rescues food and delivers it to shelters, pantries, and community partners.",
            "Our goal: turn a decade of rescue logs into actionable strategy — and a prototype coordination platform.",
        ],
    )

    add_bullet_slide(
        prs,
        "Data & scope",
        [
            "Rescue tracking dataset: donor, recipient, pounds, ZIP, dates, completion status.",
            "~155K rows across four metros; we analyze completed rescues for impact metrics.",
            "Maps use ZIP centroids; scheduling demo stores posts locally (JSON) for the hackathon.",
        ],
    )

    add_image_slide(
        prs,
        "Where volume is growing — by city (monthly lbs)",
        FIG / "01_monthly_pounds_by_city.png",
        "Finished rescues — use for growth narrative and city mix.",
    )

    add_image_slide(
        prs,
        "Seasonality — when to staff drivers",
        FIG / "02_seasonality_index.png",
        "Index vs average month — align volunteer shifts and donor reminders.",
    )

    add_image_slide(
        prs,
        "Donor types — who supplies the pounds",
        FIG / "03_top_donor_types.png",
        "Prioritize outreach by high-volume donor categories.",
    )

    add_image_slide(
        prs,
        "Top donors & recipients",
        FIG / "04_top_donors_recipients.png",
        "Deepen anchor partnerships; watch recipient load.",
    )

    add_image_slide(
        prs,
        "Operations — weekday peaks",
        FIG / "05_dow_peaks_by_city.png",
        "Schedule capacity where weekday demand spikes.",
    )

    add_image_slide(
        prs,
        "Cost focus — same-area batching (ZIP3 proxy)",
        FIG / "06_batching_share_by_city.png",
        "Multi-stop routes in dense ZIP3s reduce duplicate trip starts.",
    )

    add_bullet_slide(
        prs,
        "Maps & hot spots (live demo)",
        [
            "Route map: donor ZIP → recipient ZIP flows, weighted by pounds.",
            "Heat map: surplus concentration by donor ZIP for targeting and batching.",
            "Run locally: streamlit run app.py → Map & Waste pages.",
        ],
    )

    add_bullet_slide(
        prs,
        "Scheduling & marketplace prototype",
        [
            "Donors post pounds + pickup window; see estimated landfill savings vs trip cost.",
            "Drivers/riders post capacity and availability — live board with masked contacts.",
            "Shows how RLC could coordinate like a platform while staying nonprofit-led.",
        ],
    )

    add_bullet_slide(
        prs,
        "Route Intelligence (product modules)",
        [
            "Route batching engine: same-day + same-ZIP3 clusters → one routed run vs many solo trips; $ savings proxy.",
            "Restaurant value dashboard: landfill cost avoided vs RLC fee — net savings framing for partners.",
            "Rescue risk flagging: historical incomplete rates by weekday & ZIP3; simulator + ops queue from open rescues.",
        ],
    )

    add_bullet_slide(
        prs,
        "Economics story (for businesses)",
        [
            "Landfill/hauling is not free — donation + batched rescue can beat throw-away cost under realistic assumptions.",
            "RLC coordination spreads logistics learning; grants/partners can subsidize trips so donors are not “priced out.”",
            "Full calculator in app: Strategy & economics page.",
        ],
    )

    add_bullet_slide(
        prs,
        "Recommendations",
        [
            "Staff and route-batch around weekday + seasonal peaks.",
            "Prioritize corridors (ZIP3) with repeated multi-pickup days.",
            "Pilot scheduling board with RLC ops; add auth, DB, and SMS/email next.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank you"
    sub = slide.placeholders[1]
    sub.text = "Questions?\n\nDemo: http://127.0.0.1:8501 (local Streamlit)"
    _style_title(slide.shapes.title)
    for p in sub.text_frame.paragraphs:
        p.font.size = Pt(22)
        p.font.color.rgb = MUTED
        p.alignment = PP_ALIGN.CENTER

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
